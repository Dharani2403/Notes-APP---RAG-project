import os
import json
import docx
import PyPDF2
import pptx
import pandas as pd
from ebooklib import epub
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image

# Optional: Install via pip if missing
# pip install tiktoken nltk langchain

import tiktoken
from langchain.text_splitter import RecursiveCharacterTextSplitter

# ---------------------------
# TEXT EXTRACTION FUNCTIONS
# ---------------------------

def extract_text_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text += f"\n[Page {page_num}]\n" + page_text + "\n"
    return text.strip()

def extract_text_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            text += f"\n[Slide {i}]\n" + "\n".join(slide_text) + "\n"
    return text.strip()

def extract_text_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()

def extract_text_epub(file_path):
    book = epub.read_epub(file_path)
    text = ""
    for item in book.get_items():
        if item.get_type() == 9:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text += soup.get_text() + "\n"
    return text.strip()

def extract_text_xlsx(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text_image(file_path):
    try:
        img = Image.open(file_path)
        img = img.convert("RGB")
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        return f"Error processing image: {e}"

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_pdf(file_path)
    elif ext == ".docx":
        return extract_text_docx(file_path)
    elif ext == ".pptx":
        return extract_text_pptx(file_path)
    elif ext == ".txt":
        return extract_text_txt(file_path)
    elif ext == ".epub":
        return extract_text_epub(file_path)
    elif ext in [".xlsx", ".xls"]:
        return extract_text_xlsx(file_path)
    elif ext in [".png",".jpeg",".jpg",".webp",".bmp"]:
        return extract_text_image(file_path)
    else:
        return None

# ---------------------------
# CHUNKING FUNCTIONS
# ---------------------------

def split_text_tokens(text, chunk_size=500, overlap=50, model="gpt-3.5-turbo"):
    """Token-based chunking with tiktoken"""
    enc = tiktoken.encoding_for_model(model)
    tokens = enc.encode(text)

    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunk_tokens = tokens[start:end]
        chunk_text = enc.decode(chunk_tokens).strip()
        chunks.append(chunk_text)
        start += chunk_size - overlap
    return chunks

def semantic_split_text(text, chunk_size=500, overlap=50):
    """Semantic-aware splitting with recursive strategy"""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)

# ---------------------------
# PROCESSING & STORAGE
# ---------------------------

def process_file(file_path, output_json="data/data_chunks.json", chunk_size=500, overlap=50, use_tokens=True):
    if not os.path.isfile(file_path):
        return {"status": "error", "message": "Invalid file path!"}

    os.makedirs(os.path.dirname(output_json), exist_ok=True)

    # Load existing data
    if os.path.exists(output_json):
        with open(output_json, "r", encoding="utf-8") as f:
            try:
                data = json.load(f)
            except:
                data = []
    else:
        data = []

    # Extract text
    text = extract_text(file_path)
    if not text:
        return {"status": "error", "message": "No text extracted!"}

    # Split text into chunks
    if use_tokens:
        chunks = split_text_tokens(text, chunk_size, overlap)
    else:
        chunks = semantic_split_text(text, chunk_size, overlap)

    # Add chunks with metadata
    for i, chunk in enumerate(chunks, 1):
        entry = {
            "file_name": os.path.basename(file_path),
            "chunk_id": i,
            "content": chunk
        }
        data.append(entry)

    # Overwrite JSON
    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)

    return {"status": "success", "message": f"Extracted {len(chunks)} chunks from {file_path}"}

def run_file_extract(file_path):
    return process_file(file_path)

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        print(run_file_extract(sys.argv[1]))
